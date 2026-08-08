"""
PPT generation shared library — extracts helpers from create_ppt_v4.py.
Usage: from ppt_lib import *
"""
import os, json
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import lxml.etree as ET

# -----------------------------------------------------------------------
# Color palette
# -----------------------------------------------------------------------
RED   = RGBColor(205, 32, 38)
GOLD  = RGBColor(212, 175, 55)
WHITE = RGBColor(255, 255, 255)
OFF_W = RGBColor(235, 235, 240)
LTGY  = RGBColor(200, 200, 210)
MDGY  = RGBColor(160, 160, 170)
DKBG  = RGBColor(8, 8, 35)
NAVY  = RGBColor(15, 15, 50)
BLUE  = RGBColor(40, 60, 130)

# -----------------------------------------------------------------------
# Layout constants
# -----------------------------------------------------------------------
SW = Emu(12191365)
SH = Emu(6858000)
LI = 6                          # blank slide layout index
NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# -----------------------------------------------------------------------
# Slide image mapping (key -> fallback key) and per-slide overlay opacity
# -----------------------------------------------------------------------
SLIDE_IMGS = [
    ["republican_rally", "us_capitol"],
    ["us_capitol", "white_house"],
    ["statue_of_liberty", "grand_canyon"],
    ["mount_rushmore", "new_york_city"],
    ["statue_of_liberty", "liberty_bell"],
    ["us_capitol", "times_square"],
    ["times_square", "new_york_city"],
    ["new_york_city", "times_square"],
    ["mount_rushmore", "republican_rally"],
    ["american_flag_eagle", "white_house"],
]

OPACITIES = [
    "12000000", "05000000", "12000000", "12000000", "12000000",
    "05000000", "12000000", "12000000", "08000000", "12000000",
]

# -----------------------------------------------------------------------
# Low-level XML helper — set fill alpha
# -----------------------------------------------------------------------
def ap(sh, av):
    sp = sh._element.spPr
    sf = sp.find(NS + "solidFill")
    if sf is not None:
        sc = sf.find(NS + "srgbClr")
        if sc is not None:
            a = ET.SubElement(sc, NS + "alpha")
            a.set("val", av)

# -----------------------------------------------------------------------
# Background: picture + dark overlay at given opacity
# -----------------------------------------------------------------------
def BG(sl, img_dir, idx):
    keys = SLIDE_IMGS[idx]
    op = OPACITIES[idx]
    for k in keys:
        fp = os.path.join(img_dir, k + ".jpg")
        if os.path.exists(fp):
            sl.shapes.add_picture(fp, 0, 0, SW, SH)
            ov = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
            ov.fill.solid(); ov.fill.fore_color.rgb = DKBG
            ap(ov, op); ov.line.fill.background()
            return

# -----------------------------------------------------------------------
# Text box (single-line helper)
# -----------------------------------------------------------------------
def Tb(sl, l, t, w, h, tx, sz=18, b=False, c=WHITE, a=PP_ALIGN.LEFT):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tx
    p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = c
    p.font.name = "Microsoft YaHei"; p.alignment = a
    return bx

# -----------------------------------------------------------------------
# Accent bar (thin coloured rectangle)
# -----------------------------------------------------------------------
def Ac(sl, x, y, w, c=RED, h=Emu(5000)):
    r = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = c; r.line.fill.background()

# -----------------------------------------------------------------------
# Bottom / top coloured strips
# -----------------------------------------------------------------------
def Btm(sl): Ac(sl, 0, int(SH - 7000), SW, RED, Emu(7000))
def Top(sl): Ac(sl, 0, 0, SW, RED, Emu(7000))

# -----------------------------------------------------------------------
# Slide header block (Chinese + English title + accent bar)
# -----------------------------------------------------------------------
def Hd(sl, cn, en):
    Tb(sl, Emu(600000), Emu(350000), Emu(int(SW - 1200000)), Emu(480000),
        cn, 36, True, WHITE)
    Tb(sl, Emu(600000), Emu(850000), Emu(int(SW - 1200000)), Emu(280000),
        en, 14, False, GOLD)
    Ac(sl, Emu(600000), Emu(1160000), Emu(900000), RED, Emu(5000))

# -----------------------------------------------------------------------
# Semi-transparent rounded panel
# -----------------------------------------------------------------------
def Pn(sl, x, y, w, h, bg=None, al=None):
    p = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    g = bg if bg else DKBG
    a = al if al else "70000000"
    p.fill.solid(); p.fill.fore_color.rgb = g; ap(p, a)
    p.line.color.rgb = RED; p.line.width = Pt(1.5)

# -----------------------------------------------------------------------
# Card: rounded rectangle with coloured top accent
# -----------------------------------------------------------------------
def Cd(sl, x, y, w, h, co=RED):
    c = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = DKBG; ap(c, "70000000")
    c.line.color.rgb = co; c.line.width = Pt(1.5)
    Ac(sl, x, y, w, co, Emu(5000))
    return c

# -----------------------------------------------------------------------
# Presentation factory
# -----------------------------------------------------------------------
def new_presentation():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs

# -----------------------------------------------------------------------
# Load texts.json dict
# -----------------------------------------------------------------------
def load_texts(path):
    with open(path, "r", encoding="utf-8") as fh:
        return json.load(fh)
