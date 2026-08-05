import sys
sys.path.append(r"C:\Users\FYH\AppData\Local\Programs\Python\Python313\Lib\site-packages")
from pptx import Presentation

prs = Presentation(r"C:\Users\FYH\Desktop\banzhang_jingxuan_fangzhenping_new.pptx")

with open(r"C:\Users\FYH\Documents\New project\ppt_report.txt", "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"=== Slide {i+1} ===\n")
        for shape in slide.shapes:
            has_img = shape.shape_type == 13
            has_txt = shape.has_text_frame
            txt = ""
            if has_txt:
                for p in shape.text_frame.paragraphs:
                    txt += p.text
            if has_img:
                f.write(f"  [IMAGE]\n")
            elif txt.strip():
                f.write(f'  [TEXT] {txt.strip()[:80]}\n')
        f.write("\n")

import os
size = os.path.getsize(r"C:\Users\FYH\Desktop\banzhang_jingxuan_fangzhenping_new.pptx")
f.write(f"File size: {size} bytes ({size/1024/1024:.1f} MB)\n")
print("Report written successfully")
