import sys, os
sys.stdout.reconfigure(encoding='utf-8')
sys.path.append(r'C:\Users\FYH\AppData\Local\Programs\Python\Python313\Lib\site-packages')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

pptx_path = 'C:/Users/FYH/Desktop/\u73ed\u957f\u7ade\u9009_\u623f\u632f\u5e73.pptx'
prs = Presentation(pptx_path)
print(f'Slide width: {prs.slide_width} EMU')
print(f'Slide height: {prs.slide_height} EMU')
print(f'Slide count: {len(prs.slides)}')
print()
for i, slide in enumerate(prs.slides):
    print(f'=== Slide {i+1} ===')
    print(f'Layout: {slide.slide_layout.name}')
    bg = slide.background
    print(f'Background fill: {bg.fill.type}')
    for shape in slide.shapes:
        print(f'  Shape: {shape.shape_type}, name="{shape.name}"')
        print(f'    Position: left={shape.left}, top={shape.top}')
        print(f'    Size: width={shape.width}, height={shape.height}')
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font = run.font
                    c = font.color.rgb if font.color and font.color.type else None
                    txt = run.text.replace('\n', '\\n').replace('\r', '\\r')
                    print(f'    Text: "{txt}"')
                    print(f'      Font: size={font.size}, bold={font.bold}, color={c}')
                    print(f'      Alignment: {para.alignment}')
    print()
