import sys, os
sys.path.append(r"C:\Users\FYH\AppData\Local\Programs\Python\Python313\Lib\site-packages")
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import lxml.etree as ET

img_dir = r"C:\Users\FYH\Documents\New project\nordic_imgs"
slide_w = Emu(12191365)
slide_h = Emu(6858000)

def alpha_val(pct):
    """Convert 0-1 opacity to PPTX alpha value (1/1000ths of percent)"""
    return str(int(pct * 100000))

def add_bg(slide, img_path, darkness=0.3):
    """Add background image with subtle dark overlay. darkness=0.3 means 30% black overlay."""
    slide.shapes.add_picture(img_path, 0, 0, slide_w, slide_h)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    spPr = overlay._element.spPr
    sf = spPr.find(f"{ns}solidFill")
    if sf is not None:
        sc = sf.find(f"{ns}srgbClr")
        if sc is not None:
            a = ET.SubElement(sc, f"{ns}alpha")
            a.set("val", alpha_val(darkness))
    overlay.line.fill.background()

def tb(slide, l, t, w, h, text, sz=18, b=False, c=RGBColor(255,255,255), a=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = b
    p.font.color.rgb = c
    p.font.name = "Microsoft YaHei"
    p.alignment = a

def accent(slide, l, t, w, clr=RGBColor(212,104,74)):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Emu(4000))
    r.fill.solid()
    r.fill.fore_color.rgb = clr
    r.line.fill.background()

def overlay(slide, l, t, w, h, pct=0.6, clr=RGBColor(0,0,0)):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = clr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    spPr = r._element.spPr
    sf = spPr.find(f"{ns}solidFill")
    if sf is not None:
        sc = sf.find(f"{ns}srgbClr")
        if sc is not None:
            a = ET.SubElement(sc, f"{ns}alpha")
            a.set("val", alpha_val(pct))
    r.line.fill.background()
    return r

prs = Presentation()
prs.slide_width = slide_w
prs.slide_height = slide_h
bl = prs.slide_layouts[6]

# ============= SLIDE 1: Title =============
s = prs.slides.add_slide(bl)
add_bg(s, os.path.join(img_dir, "img1.jpg"), 0.3)
accent(s, 0, 0, slide_w)
cx, cw = Emu(int(slide_w*0.12)), Emu(int(slide_w*0.76))
tb(s, cx, Emu(1800000), cw, Emu(400000), "2025级 材料成型及控制工程一班", 14, False, RGBColor(212,104,74), PP_ALIGN.CENTER)
tb(s, cx, Emu(2300000), cw, Emu(800000), "班长竞选", 52, True, RGBColor(255,255,255), PP_ALIGN.CENTER)
tb(s, cx, Emu(3200000), cw, Emu(500000), "房 振 平", 32, True, RGBColor(212,104,74), PP_ALIGN.CENTER)
tb(s, cx, Emu(3800000), cw, Emu(350000), "用心服务  \u00b7  用行担当", 16, False, RGBColor(200,200,200), PP_ALIGN.CENTER)
accent(s, 0, int(slide_h-6000), slide_w)

# ============= SLIDE 2: About Me =============
s = prs.slides.add_slide(bl)
add_bg(s, os.path.join(img_dir, "img2.jpg"), 0.3)
pw = Emu(int(slide_w*0.35))
overlay(s, 0, 0, pw, slide_h, 0.65, RGBColor(8,8,20))

tb(s, Emu(350000), Emu(500000), Emu(pw-700000), Emu(400000), "关于我", 32, True, RGBColor(255,255,255))
accent(s, Emu(350000), Emu(950000), Emu(700000))
tb(s, Emu(350000), Emu(1050000), Emu(pw-700000), Emu(300000), "有经验  \u00b7  有热情  \u00b7  有担当", 13, False, RGBColor(180,180,180))

ph = Emu(1500000)
pc = s.shapes.add_shape(MSO_SHAPE.OVAL, int((pw-ph)/2), Emu(1700000), ph, ph)
pc.fill.solid()
pc.fill.fore_color.rgb = RGBColor(50,50,70)
pc.line.color.rgb = RGBColor(212,104,74)
pc.line.width = Pt(3)

tb(s, Emu(100000), Emu(3400000), Emu(pw-200000), Emu(350000), "房振平", 22, True, RGBColor(255,255,255), PP_ALIGN.CENTER)
tb(s, Emu(100000), Emu(3750000), Emu(pw-200000), Emu(250000), "大二材料成型及控制工程一班", 11, False, RGBColor(160,160,160), PP_ALIGN.CENTER)

rl, rw = Emu(int(slide_w*0.38)), Emu(int(slide_w*0.57))
tb(s, rl, Emu(400000), rw, Emu(400000), "我的经历", 28, True, RGBColor(255,255,255))
accent(s, rl, Emu(850000), Emu(500000))

tlx = rl + Emu(250000)
tl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tlx, Emu(1200000), Emu(4000), Emu(4500000))
tl.fill.solid()
tl.fill.fore_color.rgb = RGBColor(212,104,74)
tl.line.fill.background()

d2 = [
    ("初中 \u00b7 班长（三年）", "连续三年担任班长，积累了丰富的\n班级管理经验，善于组织协调与沟通", RGBColor(212,104,74)),
    ("高中 \u00b7 团支书", "负责团支部工作，搭建师生沟通桥梁，\n组织团日活动，提升班级凝聚力", RGBColor(240,160,80)),
    ("现在 \u00b7 竞选班长", "以过去的经验为基础，用心服务材控一班，\n用行担当不负大家期盼", RGBColor(100,160,120)),
]
for i, (t, d, clr) in enumerate(d2):
    yb = Emu(1100000 + i*1600000)
    ds = s.shapes.add_shape(MSO_SHAPE.OVAL, tlx-Emu(80000), yb, Emu(180000), Emu(180000))
    ds.fill.solid()
    ds.fill.fore_color.rgb = clr
    ds.line.fill.background()
    tb(s, rl+Emu(350000), yb-Emu(50000), rw-Emu(350000), Emu(300000), t, 18, True, clr)
    tb(s, rl+Emu(350000), yb+Emu(300000), rw-Emu(350000), Emu(400000), d, 12, False, RGBColor(200,200,200))

# ============= SLIDE 3: Strengths =============
s = prs.slides.add_slide(bl)
add_bg(s, os.path.join(img_dir, "img3.jpg"), 0.28)

tb(s, Emu(600000), Emu(300000), Emu(slide_w-1200000), Emu(400000), "我的优势", 32, True, RGBColor(255,255,255))
accent(s, Emu(600000), Emu(750000), Emu(800000))
tb(s, Emu(600000), Emu(850000), Emu(slide_w-1200000), Emu(250000), "为什么我可以胜任班长一职？", 14, False, RGBColor(180,180,180))

cards = [
    ("\U0001F3C6", "组织经验丰富", "初中三年班长 + 高中团支书\n深谙班级管理之道", RGBColor(212,104,74)),
    ("\U0001F91D", "沟通能力出众", "善于倾听同学心声\n搭建师生沟通桥梁", RGBColor(240,160,80)),
    ("\u2B50", "责任心强", "做事有始有终\n用心对待每一件事、每一个人", RGBColor(100,160,120)),
    ("\U0001F680", "执行力到位", "策划组织多次班级活动\n让计划落地，让想法成真", RGBColor(235,150,120)),
]
cw = Emu(int((slide_w-1800000)/2))
ch = Emu(2200000)
g = Emu(250000)
sx, sy = Emu(600000), Emu(1400000)

for i, (icon, t, d, clr) in enumerate(cards):
    col, row = i%2, i//2
    x, y = sx+col*(cw+g), sy+row*(ch+g)
    cd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
    cd.fill.solid()
    cd.fill.fore_color.rgb = RGBColor(12,12,30)
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    spPr = cd._element.spPr
    sf = spPr.find(f"{ns}solidFill")
    if sf is not None:
        sc = sf.find(f"{ns}srgbClr")
        if sc is not None:
            a = ET.SubElement(sc, f"{ns}alpha")
            a.set("val", alpha_val(0.6))
    cd.line.color.rgb = clr
    cd.line.width = Pt(1.5)
    accent(s, x, y, cw, clr)
    tb(s, x+Emu(150000), y+Emu(150000), Emu(600000), Emu(450000), icon, 28, False, RGBColor(255,255,255))
    tb(s, x+Emu(700000), y+Emu(180000), Emu(cw-900000), Emu(350000), t, 18, True, clr)
    tb(s, x+Emu(150000), y+Emu(700000), Emu(cw-300000), Emu(ch-900000), d, 12, False, RGBColor(210,210,210))

# ============= SLIDE 4: Platform =============
s = prs.slides.add_slide(bl)
add_bg(s, os.path.join(img_dir, "img4.jpg"), 0.3)

tb(s, Emu(600000), Emu(300000), Emu(slide_w-1200000), Emu(400000), "施政纲领", 32, True, RGBColor(255,255,255))
accent(s, Emu(600000), Emu(750000), Emu(800000))
tb(s, Emu(600000), Emu(850000), Emu(slide_w-1200000), Emu(250000), "如果当选，我承诺用心做好以下四件事", 14, False, RGBColor(180,180,180))

platform = [
    ("01", "沟通桥梁", "定期收集同学意见\n及时反馈给老师\n让班级沟通零距离", RGBColor(212,104,74)),
    ("02", "班级活动", "每月组织有温度的班级活动\n增强凝聚力\n丰富课余生活", RGBColor(240,160,80)),
    ("03", "学习氛围", "组建学习互助小组\n共享资源\n带动大家一起进步", RGBColor(100,160,120)),
    ("04", "公平透明", "班级事务公开透明\n尊重每一个人的声音\n公平公正", RGBColor(235,150,120)),
]
iw = Emu(int((slide_w-2000000)/4))
ig = Emu(200000)
sx2 = Emu(600000)

for i, (num, t, d, clr) in enumerate(platform):
    x = sx2 + i*(iw+ig)
    y = Emu(1600000)
    cs = Emu(450000)
    cx = int(x+(iw-cs)/2)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, y, cs, cs)
    c.fill.solid()
    c.fill.fore_color.rgb = clr
    c.line.fill.background()
    tb(s, cx, y, cs, cs, num, 20, True, RGBColor(255,255,255), PP_ALIGN.CENTER)
    tb(s, x, y+Emu(550000), iw, Emu(350000), t, 18, True, clr, PP_ALIGN.CENTER)
    tb(s, x+Emu(30000), y+Emu(950000), Emu(iw-60000), Emu(500000), d, 12, False, RGBColor(210,210,210), PP_ALIGN.CENTER)

# Bottom quote
overlay(s, 0, int(slide_h-650000), slide_w, Emu(650000), 0.6)
tb(s, Emu(int(slide_w*0.1)), int(slide_h-550000), Emu(int(slide_w*0.8)), Emu(350000),
   "\u300C用心服务材控一班，用行担当不负期盼\u300D", 16, False, RGBColor(212,104,74), PP_ALIGN.CENTER)

# ============= SLIDE 5: Closing =============
s = prs.slides.add_slide(bl)
add_bg(s, os.path.join(img_dir, "img5.jpg"), 0.3)
accent(s, 0, 0, slide_w)

tb(s, Emu(1500000), Emu(1800000), Emu(slide_w-3000000), Emu(500000), "用心服务  \u00b7  用行担当", 36, True, RGBColor(255,255,255), PP_ALIGN.CENTER)
tb(s, Emu(1500000), Emu(2500000), Emu(slide_w-3000000), Emu(400000), "请投我一票", 28, False, RGBColor(212,104,74), PP_ALIGN.CENTER)
tb(s, Emu(1500000), Emu(3000000), Emu(slide_w-3000000), Emu(350000), "\u2014\u2014 房振平", 20, False, RGBColor(200,200,200), PP_ALIGN.CENTER)
tb(s, Emu(1500000), Emu(3450000), Emu(slide_w-3000000), Emu(250000), "大二材料成型及控制工程一班", 13, False, RGBColor(160,160,160), PP_ALIGN.CENTER)
accent(s, int(slide_w/2-600000), Emu(4000000), Emu(1200000))
tb(s, Emu(2000000), Emu(4200000), Emu(slide_w-4000000), Emu(250000), "你的每一票，都是我前行的动力", 14, False, RGBColor(180,180,180), PP_ALIGN.CENTER)
accent(s, 0, int(slide_h-6000), slide_w)

out = r"C:\Users\FYH\Desktop\班长竞选_房振平_新版.pptx"
prs.save(out)
print(f"OK: {out}")
