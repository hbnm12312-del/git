import json, sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import lxml.etree as ET
T = json.load(open(r"C:\Users\FYH\Documents\New project\texts.json", "r", encoding="utf-8"))
IMG = r"C:\Users\FYH\Documents\New project\us_imgs_hd"
OUT = r"C:\Users\FYH\Documents\New project\Ding_Dongxu_President.pptx"
SLIDE_IMGS = [
    ["republican_rally", "us_capitol"],
    ["us_capitol", "white_house"],
    ["statue_of_liberty", "grand_canyon"],
    ["mount_rushmore", "new_york_city"],
    ["statue_of_liberty", "liberty_bell"],
    ["us_capitol", "times_square"],
    ["times_square", "new_york_city"],
    ["new_york_city", "times_square"],
    ["mount_rushmore", "republican_rally"],
    ["american_flag_eagle", "white_house"],
]
OPS = ["12000000","05000000","12000000","12000000","12000000",
       "05000000","12000000","12000000","08000000","12000000"]
RED=RGBColor(205,32,38);GOLD=RGBColor(212,175,55);WHITE=RGBColor(255,255,255)
OFF_W=RGBColor(235,235,240);LTGY=RGBColor(200,200,210);MDGY=RGBColor(160,160,170)
DKBG=RGBColor(8,8,35);NAVY=RGBColor(15,15,50);BLUE=RGBColor(40,60,130)
SW=Emu(12191365);SH=Emu(6858000);LI=6
NS="{http://schemas.openxmlformats.org/drawingml/2006/main}"
def ap(sh,av):
    sp=sh._element.spPr;sf=sp.find(NS+"solidFill")
    if sf is not None:
        sc=sf.find(NS+"srgbClr")
        if sc is not None:
            a=ET.SubElement(sc,NS+"alpha");a.set("val",av)
def BG(sl,idx):
    keys=SLIDE_IMGS[idx];op=OPS[idx]
    for k in keys:
        fp=os.path.join(IMG,k+".jpg")
        if os.path.exists(fp):
            sl.shapes.add_picture(fp,0,0,SW,SH)
            ov=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
            ov.fill.solid();ov.fill.fore_color.rgb=DKBG;ap(ov,op)
            ov.line.fill.background();return
def Tb(sl,l,t,w,h,tx,sz=18,b=False,c=WHITE,a=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(l,t,w,h)
    tf=bx.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=tx
    p.font.size=Pt(sz);p.font.bold=b;p.font.color.rgb=c
    p.font.name="Microsoft YaHei";p.alignment=a;return bx
def Ac(sl,x,y,w,c=RED,h=Emu(5000)):
    r=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    r.fill.solid();r.fill.fore_color.rgb=c;r.line.fill.background()
def Btm(sl): Ac(sl,0,int(SH-7000),SW,RED,Emu(7000))
def Top(sl): Ac(sl,0,0,SW,RED,Emu(7000))
def Hd(sl,cn,en):
    Tb(sl,Emu(600000),Emu(350000),Emu(int(SW-1200000)),Emu(480000),cn,36,True,WHITE)
    Tb(sl,Emu(600000),Emu(850000),Emu(int(SW-1200000)),Emu(280000),en,14,False,GOLD)
    Ac(sl,Emu(600000),Emu(1160000),Emu(900000),RED,Emu(5000))
def Pn(sl,x,y,w,h,bg=None,al=None):
    p=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h)
    g=bg if bg else DKBG; a=al if al else "70000000"
    p.fill.solid();p.fill.fore_color.rgb=g;ap(p,a)
    p.line.color.rgb=RED;p.line.width=Pt(1.5)
def Cd(sl,x,y,w,h,co=RED):
    c=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h)
    c.fill.solid();c.fill.fore_color.rgb=DKBG;ap(c,"70000000")
    c.line.color.rgb=co;c.line.width=Pt(1.5)
    Ac(sl,x,y,w,co,Emu(5000));return c
prs=Presentation();prs.slide_width=SW;prs.slide_height=SH
print("Start")
# === SLIDE 1 ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
BG(s, 0); Top(s); Btm(s)
Ac(s, 0, Emu(100000), SW, GOLD, Emu(2000))
Ac(s, 0, int(SH-100000), SW, GOLD, Emu(2000))
Tb(s, Emu(800000), Emu(800000), Emu(int(SW-1600000)), Emu(250000), T["s1_gop"], 11, False, GOLD, PP_ALIGN.CENTER)
Tb(s, Emu(800000), Emu(1200000), Emu(int(SW-1600000)), Emu(700000), T["s1_name_cn"], 56, True, WHITE, PP_ALIGN.CENTER)
Tb(s, Emu(800000), Emu(1900000), Emu(int(SW-1600000)), Emu(400000), T["s1_name_en"], 40, True, WHITE, PP_ALIGN.CENTER)
# User photo on title slide
up_path = os.path.join(IMG, "user_photo_slide1.jpg")
if os.path.exists(up_path):
    try:
        circ = Emu(1200000)
        cx = int(SW/2 - circ/2)
        s.shapes.add_picture(up_path, cx, Emu(2950000), circ, circ)
        border = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, Emu(2950000), circ, circ)
        border.fill.background(); border.line.color.rgb = GOLD; border.line.width = Pt(3)
    except:
        pass
Tb(s, Emu(800000), Emu(2500000), Emu(int(SW-1600000)), Emu(350000), T["s1_slogan"], 18, True, GOLD, PP_ALIGN.CENTER)
Tb(s, Emu(800000), Emu(2950000), Emu(int(SW-1600000)), Emu(280000), T["s1_sub"], 14, False, OFF_W, PP_ALIGN.CENTER)
st = "\u2605 \u00a0 \u2605 \u00a0 \u2605 \u00a0 \u2605 \u00a0 \u2605"
Tb(s, Emu(800000), Emu(3800000), Emu(int(SW-1600000)), Emu(200000), st, 16, False, GOLD, PP_ALIGN.CENTER)
Tb(s, Emu(800000), Emu(4100000), Emu(int(SW-1600000)), Emu(250000), T["s1_career_cn"], 13, False, LTGY, PP_ALIGN.CENTER)
Tb(s, Emu(800000), Emu(4350000), Emu(int(SW-1600000)), Emu(250000), T["s1_career_en"], 11, False, MDGY, PP_ALIGN.CENTER)
Cd(s, Emu(int(SW/2-2500000)), Emu(5100000), Emu(5000000), Emu(500000))
Tb(s, Emu(int(SW/2-2300000)), Emu(5150000), Emu(4600000), Emu(400000), T["s1_bottom"], 16, True, GOLD, PP_ALIGN.CENTER)
print("Slide 1 done")

# === SLIDE 2: About Me ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
BG(s, 1)
Pn(s, Emu(400000), Emu(300000), Emu(4000000), int(SH-800000), NAVY, "80000000")
Tb(s, Emu(500000), Emu(500000), Emu(3800000), Emu(500000), T["s2_title_cn"], 36, True, WHITE, PP_ALIGN.CENTER)
Tb(s, Emu(500000), Emu(950000), Emu(3800000), Emu(300000), T["s2_title_en"], 16, False, GOLD, PP_ALIGN.CENTER)
# Profile photo - use user photo
user_photo = os.path.join(IMG, "285d1b8cc9e3bb597bb381e7e545907.jpg")
if os.path.exists(user_photo):
    circ = Emu(1700000); cx = int(Emu(500000)+(Emu(3800000)-circ)/2)
    try:
        s.shapes.add_picture(user_photo, cx, Emu(1500000), circ, circ)
        # Add oval clipping via a circle overlay with white border
        border = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, Emu(1500000), circ, circ)
        border.fill.background(); border.line.color.rgb = RED; border.line.width = Pt(4)
    except:
        pass
else:
    circ = Emu(1700000); cx = int(Emu(500000)+(Emu(3800000)-circ)/2)
    ci = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, Emu(1500000), circ, circ)
    ci.fill.solid(); ci.fill.fore_color.rgb = RGBColor(30,30,60); ci.line.color.rgb = RED; ci.line.width = Pt(4)
    Tb(s, cx, Emu(2000000), circ, Emu(400000), T["s2_name"], 26, True, WHITE, PP_ALIGN.CENTER)
Tb(s, cx, Emu(2400000), circ, Emu(280000), "Ding Dongxu", 14, False, GOLD, PP_ALIGN.CENTER)
Tb(s, cx, Emu(2680000), circ, Emu(200000), "1970 -", 12, False, LTGY, PP_ALIGN.CENTER)
Tb(s, Emu(500000), Emu(3650000), Emu(3800000), Emu(200000), T["s2_edu_title"], 13, True, GOLD, PP_ALIGN.CENTER)
edu = [T["s2_edu_1"], T["s2_edu_2"], T["s2_edu_3"]]
for i, e in enumerate(edu):
    Tb(s, Emu(500000), Emu(3900000+i*250000), Emu(3800000), Emu(220000), "\u25b8  " + e, 11, False, OFF_W)
# Right side - career timeline
rx = Emu(4800000); rw = Emu(int(SW-5200000))
Tb(s, rx, Emu(500000), rw, Emu(400000), T["s2_career_title_cn"], 30, True, WHITE)
Tb(s, rx, Emu(880000), rw, Emu(250000), T["s2_career_title_en"], 14, False, GOLD)
Ac(s, rx, Emu(1180000), Emu(600000), RED, Emu(5000))
tlx = rx + Emu(150000)
tl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tlx, Emu(1450000), Emu(4000), Emu(4700000))
tl.fill.solid(); tl.fill.fore_color.rgb = RED; tl.line.fill.background()
career_data = [("s2_c1", T["s2_c1_per"], T["s2_c1_tit"], T["s2_c1_desc"]),
               ("s2_c2", T["s2_c2_per"], T["s2_c2_tit"], T["s2_c2_desc"]),
               ("s2_c3", T["s2_c3_per"], T["s2_c3_tit"], T["s2_c3_desc"]),
               ("s2_c4", T["s2_c4_per"], T["s2_c4_tit"], T["s2_c4_desc"])]
for i, (k, per, tit, desc) in enumerate(career_data):
    yb = Emu(1450000+i*1200000)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, tlx-Emu(60000), yb, Emu(130000), Emu(130000))
    dot.fill.solid(); dot.fill.fore_color.rgb = GOLD; dot.line.fill.background()
    Tb(s, tlx+Emu(250000), yb-Emu(50000), rw-Emu(250000), Emu(250000), per, 10, False, GOLD)
    Tb(s, tlx+Emu(250000), yb+Emu(200000), rw-Emu(250000), Emu(280000), tit, 18, True, WHITE)
    Tb(s, tlx+Emu(250000), yb+Emu(480000), rw-Emu(250000), Emu(400000), desc, 11, False, LTGY)
Btm(s)
print("Slide 2 done")
# === SLIDE 3: Vision ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
BG(s, 2); Btm(s)
Hd(s, T["s3_title_cn"], T["s3_title_en"])
visions = [("\u2605", T["s3_v1_cn"], T["s3_v1_en"], T["s3_v1_desc"]),
           ("\u2605", T["s3_v2_cn"], T["s3_v2_en"], T["s3_v2_desc"]),
           ("\u2605", T["s3_v3_cn"], T["s3_v3_en"], T["s3_v3_desc"])]
cw = Emu(3400000); ch = Emu(3500000); cgap = Emu(250000)
sx = Emu(int((SW-(cw*3+cgap*2))/2))
for i, (star, cn, en, desc) in enumerate(visions):
    cx = sx + i*(cw+cgap); cy = Emu(1600000)
    Cd(s, cx, cy, cw, ch)
    Tb(s, cx, cy+Emu(200000), cw, Emu(400000), star, 30, False, GOLD, PP_ALIGN.CENTER)
    Tb(s, cx, cy+Emu(700000), cw, Emu(400000), cn, 28, True, WHITE, PP_ALIGN.CENTER)
    Tb(s, cx, cy+Emu(1100000), cw, Emu(300000), en, 14, False, GOLD, PP_ALIGN.CENTER)
    Ac(s, int(cx+cw/2-300000), cy+Emu(1500000), Emu(600000), GOLD, Emu(3000))
    Tb(s, cx+Emu(300000), cy+Emu(1700000), cw-Emu(600000), Emu(1500000), desc, 14, False, LTGY, PP_ALIGN.CENTER)
Tb(s, Emu(1500000), int(SH-900000), Emu(int(SW-3000000)), Emu(350000), "\u201c" + T["s3_quote"] + "\u201d", 15, False, GOLD, PP_ALIGN.CENTER)
Tb(s, Emu(1500000), int(SH-600000), Emu(int(SW-3000000)), Emu(250000), "\u201c" + T["s3_quote_en"] + "\u201d", 12, False, LTGY, PP_ALIGN.CENTER)
print("Slide 3 done")

# === SLIDE 4: Economy ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
BG(s, 3); Btm(s)
Hd(s, T["s4_title_cn"], T["s4_title_en"])
Pn(s, Emu(400000), Emu(1500000), Emu(6200000), int(SH-2400000))
Tb(s, Emu(600000), Emu(1700000), Emu(5800000), Emu(350000), "\u261b  " + T["s4_sub"], 18, True, GOLD)
pols = [("01", T["s4_p1_cn"], T["s4_p1_en"], T["s4_p1_desc"]),
         ("02", T["s4_p2_cn"], T["s4_p2_en"], T["s4_p2_desc"]),
         ("03", T["s4_p3_cn"], T["s4_p3_en"], T["s4_p3_desc"]),
         ("04", T["s4_p4_cn"], T["s4_p4_en"], T["s4_p4_desc"])]
for i, (num, cn, en, desc) in enumerate(pols):
    yb = Emu(2200000+i*900000)
    ns = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(700000), yb, Emu(500000), Emu(500000))
    ns.fill.solid(); ns.fill.fore_color.rgb = RED; ns.line.fill.background()
    Tb(s, Emu(700000), yb+Emu(70000), Emu(500000), Emu(380000), num, 16, True, WHITE, PP_ALIGN.CENTER)
    Tb(s, Emu(1350000), yb+Emu(20000), Emu(5000000), Emu(260000), cn+"  /  "+en, 16, True, WHITE)
    Tb(s, Emu(1350000), yb+Emu(280000), Emu(5000000), Emu(250000), desc, 12, False, LTGY)
r2x = Emu(6900000); r2w = Emu(int(SW-7400000))
Pn(s, r2x, Emu(1500000), r2w, int(SH-2400000))
Tb(s, r2x+Emu(200000), Emu(1700000), r2w-Emu(400000), Emu(250000), T["s4_goal"], 16, True, GOLD, PP_ALIGN.CENTER)
Ac(s, int(r2x+r2w/2-300000), Emu(2000000), Emu(600000), RED, Emu(3000))
for i, (num, lab) in enumerate([(T["s4_s1_num"], T["s4_s1_lab"]),
    (T["s4_s2_num"], T["s4_s2_lab"]), (T["s4_s3_num"], T["s4_s3_lab"])]):
    ys = Emu(2300000+i*1200000)
    Tb(s, r2x+Emu(200000), ys, r2w-Emu(400000), Emu(350000), num, 36, True, GOLD, PP_ALIGN.CENTER)
    Tb(s, r2x+Emu(200000), ys+Emu(380000), r2w-Emu(400000), Emu(250000), lab, 13, True, WHITE, PP_ALIGN.CENTER)
print("Slide 4 done")

# === SLIDE 5: Immigration ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
BG(s, 4); Btm(s)
Hd(s, T["s5_title_cn"], T["s5_title_en"])
Pn(s, Emu(400000), Emu(1500000), int(SW-800000), int(SH-2400000))
col_w = Emu(int((SW-1800000)/2))
Tb(s, Emu(600000), Emu(1700000), col_w, Emu(300000), "\u26d4  " + T["s5_left_title"], 18, True, RED)
left = [T["s5_l1"], T["s5_l2"], T["s5_l3"], T["s5_l4"], T["s5_l5"]]
for i, it in enumerate(left):
    Tb(s, Emu(750000), Emu(2100000+i*320000), col_w+Emu(500000), Emu(280000), "\u2716  " + it, 12, False, OFF_W)
rx2 = Emu(6600000)
Tb(s, rx2, Emu(1700000), col_w, Emu(300000), "\u2714\ufe0f  " + T["s5_right_title"], 18, True, GOLD)
right = [T["s5_r1"], T["s5_r2"], T["s5_r3"], T["s5_r4"], T["s5_r5"]]
for i, it in enumerate(right):
    Tb(s, rx2+Emu(100000), Emu(2100000+i*320000), col_w+Emu(500000), Emu(280000), "\u2714  " + it, 12, False, OFF_W)
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(SW/2-2000), Emu(2100000), Emu(4000), Emu(1600000))
div.fill.solid(); div.fill.fore_color.rgb = RED; div.line.fill.background()
Tb(s, Emu(800000), int(SH-800000), Emu(int(SW-1600000)), Emu(300000), "\u201c" + T["s5_quote"] + "\u201d", 14, False, GOLD, PP_ALIGN.CENTER)
Tb(s, Emu(800000), int(SH-550000), Emu(int(SW-1600000)), Emu(250000), "\u201c" + T["s5_quote_en"] + "\u201d", 11, False, LTGY, PP_ALIGN.CENTER)
print("Slide 5 done")
# === SLIDE 6: Foreign Policy ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
BG(s, 5); Btm(s)
Hd(s, T["s6_title_cn"], T["s6_title_en"])
Pn(s, Emu(300000), Emu(1500000), int(SW-600000), int(SH-2400000))
Tb(s, Emu(500000), Emu(1700000), int(SW-1200000), Emu(300000), "\u261b  " + T["s6_sub"], 18, True, GOLD)
pillars = [("\U0001f30d", T["s6_p1_cn"], T["s6_p1_en"], T["s6_p1_desc"], T["s6_p1_en_desc"]),
           ("\U0001f30e", T["s6_p2_cn"], T["s6_p2_en"], T["s6_p2_desc"], T["s6_p2_en_desc"]),
           ("\U0001f30f", T["s6_p3_cn"], T["s6_p3_en"], T["s6_p3_desc"], T["s6_p3_en_desc"])]
pw = Emu(3600000); ph = Emu(2800000); pgap = Emu(250000)
psx = Emu(int(450000))
for i, (icon, cn, en, cnd, end) in enumerate(pillars):
    px = psx + i*(pw+pgap)
    co = RED if i==0 else (GOLD if i==1 else BLUE)
    Cd(s, px, Emu(2200000), pw, ph, co)
    Tb(s, px, Emu(2350000), pw, Emu(400000), icon, 28, False, WHITE, PP_ALIGN.CENTER)
    Tb(s, px, Emu(2750000), pw, Emu(300000), cn, 20, True, WHITE, PP_ALIGN.CENTER)
    Tb(s, px, Emu(3050000), pw, Emu(250000), en, 13, False, GOLD, PP_ALIGN.CENTER)
    Ac(s, int(px+pw/2-250000), Emu(3350000), Emu(500000), GOLD, Emu(3000))
    Tb(s, px+Emu(200000), Emu(3450000), pw-Emu(400000), Emu(1300000), cnd, 12, False, OFF_W)
    Tb(s, px+Emu(200000), Emu(4600000), pw-Emu(400000), Emu(350000), end, 10, False, MDGY)
print("Slide 6 done")

# === SLIDE 7: Infrastructure & Healthcare ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
BG(s, 6); Btm(s)
Hd(s, T["s7_title_cn"], T["s7_title_en"])
Pn(s, Emu(400000), Emu(1500000), int(SW-800000), Emu(1600000))
Tb(s, Emu(600000), Emu(1650000), Emu(5000000), Emu(300000), "\U0001f3d7\ufe0f  " + T["s7_infra_title"], 18, True, GOLD)
infra = [T["s7_infra_1"], T["s7_infra_2"], T["s7_infra_3"], T["s7_infra_4"]]
for i, it in enumerate(infra):
    Tb(s, Emu(700000), Emu(2050000+i*280000), Emu(5500000), Emu(250000), "\u25b8  " + it, 12, False, OFF_W)
rhx = Emu(6500000)
Tb(s, rhx, Emu(1650000), Emu(5000000), Emu(300000), "\U0001f3e5  " + T["s7_health_title"], 18, True, RED)
health = [T["s7_health_1"], T["s7_health_2"], T["s7_health_3"], T["s7_health_4"]]
for i, it in enumerate(health):
    Tb(s, rhx+Emu(100000), Emu(2050000+i*280000), Emu(5500000), Emu(250000), "\u25b8  " + it, 12, False, OFF_W)
div2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(SW/2-2000), Emu(2000000), Emu(4000), Emu(1200000))
div2.fill.solid(); div2.fill.fore_color.rgb = GOLD; div2.line.fill.background()
for i, (num, lab) in enumerate([(T["s7_s1_num"], T["s7_s1_lab"]),
    (T["s7_s2_num"], T["s7_s2_lab"]), (T["s7_s3_num"], T["s7_s3_lab"])]):
    sx = Emu(1200000+i*3300000)
    Tb(s, sx, Emu(3400000), Emu(2800000), Emu(350000), num, 32, True, GOLD, PP_ALIGN.CENTER)
    Tb(s, sx, Emu(3750000), Emu(2800000), Emu(250000), lab, 13, False, LTGY, PP_ALIGN.CENTER)
print("Slide 7 done")

# === SLIDE 8: Education ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
BG(s, 7); Btm(s)
Hd(s, T["s8_title_cn"], T["s8_title_en"])
Pn(s, Emu(400000), Emu(1500000), int(SW-800000), int(SH-2400000))
Tb(s, Emu(600000), Emu(1700000), int(SW-1200000), Emu(300000), "\u261b  " + T["s8_sub"], 18, True, GOLD)
reforms = [("\U0001f393", T["s8_r1_cn"], T["s8_r1_en"], T["s8_r1_desc"]),
           ("\U0001f4da", T["s8_r2_cn"], T["s8_r2_en"], T["s8_r2_desc"]),
           ("\U0001f3eb", T["s8_r3_cn"], T["s8_r3_en"], T["s8_r3_desc"]),
           ("\U0001f4bc", T["s8_r4_cn"], T["s8_r4_en"], T["s8_r4_desc"])]
for i, (icon, cn, en, desc) in enumerate(reforms):
    col = i%2; row = i//2
    ix = Emu(600000)+col*Emu(5600000); iy = Emu(2150000)+row*Emu(1450000)
    Tb(s, ix, iy, Emu(500000), Emu(400000), icon, 24, False, WHITE)
    Tb(s, ix+Emu(500000), iy, Emu(5000000), Emu(300000), cn+"  /  "+en, 16, True, WHITE)
    Tb(s, ix+Emu(500000), iy+Emu(350000), Emu(5000000), Emu(400000), desc, 11, False, OFF_W)
print("Slide 8 done")
# === SLIDE 9: Why Vote ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
BG(s, 8); Btm(s)
Hd(s, T["s9_title_cn"], T["s9_title_en"])
Pn(s, Emu(400000), Emu(1500000), int(SW-800000), int(SH-2400000))
reasons = [("\U0001f3db\ufe0f", T["s9_r1_cn"], T["s9_r1_en"], T["s9_r1_desc"]),
           ("\U0001f1fa\U0001f1f8", T["s9_r2_cn"], T["s9_r2_en"], T["s9_r2_desc"]),
           ("\U0001f4aa", T["s9_r3_cn"], T["s9_r3_en"], T["s9_r3_desc"]),
           ("\U0001f30a", T["s9_r4_cn"], T["s9_r4_en"], T["s9_r4_desc"])]
rw = Emu(2700000); rh = Emu(2100000); rgap = Emu(250000); rsx = Emu(550000)
for i, (icon, cn, en, desc) in enumerate(reasons):
    col = i%2; row = i//2
    rx = rsx+col*(rw+rgap+Emu(100000)); ry = Emu(1750000)+row*(rh+Emu(150000))
    Cd(s, rx, ry, rw, rh)
    Tb(s, rx, ry+Emu(150000), rw, Emu(350000), icon, 24, False, WHITE, PP_ALIGN.CENTER)
    Tb(s, rx, ry+Emu(500000), rw, Emu(300000), cn, 20, True, WHITE, PP_ALIGN.CENTER)
    Tb(s, rx, ry+Emu(800000), rw, Emu(250000), en, 12, False, GOLD, PP_ALIGN.CENTER)
    Ac(s, int(rx+rw/2-250000), ry+Emu(1100000), Emu(500000), GOLD, Emu(3000))
    Tb(s, rx+Emu(150000), ry+Emu(1200000), rw-Emu(300000), Emu(700000), desc, 12, False, LTGY, PP_ALIGN.CENTER)
print("Slide 9 done")

# === SLIDE 10: Call to Action ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
BG(s, 9); Top(s); Btm(s)
Ac(s, 0, Emu(100000), SW, GOLD, Emu(2000))
Ac(s, 0, int(SH-100000), SW, GOLD, Emu(2000))
Tb(s, Emu(500000), Emu(1200000), Emu(int(SW-1000000)), Emu(600000), T["s10_slogan"], 48, True, WHITE, PP_ALIGN.CENTER)
Tb(s, Emu(500000), Emu(1750000), Emu(int(SW-1000000)), Emu(350000), T["s10_slogan_en"], 36, True, GOLD, PP_ALIGN.CENTER)
st = "\u2605 \u00b7 \u2605 \u00b7 \u2605 \u00b7 \u2605 \u00b7 \u2605"
Tb(s, Emu(500000), Emu(2300000), Emu(int(SW-1000000)), Emu(300000), st, 18, False, RED, PP_ALIGN.CENTER)
Tb(s, Emu(500000), Emu(2700000), Emu(int(SW-1000000)), Emu(350000), T["s10_vote_cn"], 28, True, WHITE, PP_ALIGN.CENTER)
Tb(s, Emu(500000), Emu(3050000), Emu(int(SW-1000000)), Emu(300000), T["s10_vote_en"], 18, True, GOLD, PP_ALIGN.CENTER)
Cd(s, Emu(int(SW/2-2500000)), Emu(3600000), Emu(5000000), Emu(650000), GOLD)
Tb(s, Emu(int(SW/2-2300000)), Emu(3680000), Emu(4600000), Emu(260000), "\u201c" + T["s10_motto"] + "\u201d", 16, True, WHITE, PP_ALIGN.CENTER)
Tb(s, Emu(int(SW/2-2300000)), Emu(3950000), Emu(4600000), Emu(250000), "\u201c" + T["s10_motto_en"] + "\u201d", 12, False, GOLD, PP_ALIGN.CENTER)
Tb(s, Emu(500000), Emu(4650000), Emu(int(SW-1000000)), Emu(250000), T["s10_footer_cn"], 13, False, LTGY, PP_ALIGN.CENTER)
Tb(s, Emu(500000), Emu(4880000), Emu(int(SW-1000000)), Emu(250000), T["s10_footer_en"], 11, False, MDGY, PP_ALIGN.CENTER)
print("Slide 10 done")

# === SAVE ===
prs.save(OUT)
print("PPT generated successfully!")
print("Output: " + OUT)
print("Total slides: " + str(len(prs.slides)))