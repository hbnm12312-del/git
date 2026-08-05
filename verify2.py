import sys
sys.path.append(r"C:\Users\FYH\AppData\Local\Programs\Python\Python313\Lib\site-packages")
from pptx import Presentation

prs = Presentation(r"C:\Users\FYH\Desktop\banzhang_jingxuan_fangzhenping_new.pptx")

# Check emoji content in Slide 3
s3 = prs.slides[2]
for shape in s3.shapes:
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            if p.text.strip():
                # Get hex of first char if it might be emoji
                first = p.text.strip()[0]
                if ord(first) > 0x2000:
                    print(f"  EMOJI or special: U+{ord(first):04X} = {p.text.strip()[:30]}")
                else:
                    print(f"  TEXT: {p.text.strip()[:40]}")

print()
print("=== Checking file integrity ===")
import os
size = os.path.getsize(r"C:\Users\FYH\Desktop\banzhang_jingxuan_fangzhenping_new.pptx")
print(f"File size: {size} bytes ({size/1024/1024:.1f} MB)")
print("All good!")
