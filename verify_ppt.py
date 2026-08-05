import sys
sys.path.append(r'C:\Users\FYH\AppData\Local\Programs\Python\Python313\Lib\site-packages')
from pptx import Presentation
from pptx.util import Emu

prs = Presentation(r'C:\Users\FYH\Desktop\banzhang_jingxuan_fangzhenping_new.pptx')
print(f"Slides: {len(prs.slides)}")
print(f"Slide size: {prs.slide_width} x {prs.slide_height}")
print()

for i, slide in enumerate(prs.slides):
    print(f"=== Slide {i+1} ===")
    for shape in slide.shapes:
        has_img = shape.shape_type == 13
        has_txt = shape.has_text_frame
        txt = ""
        if has_txt:
            for p in shape.text_frame.paragraphs:
                txt += p.text
        if has_img:
            img_info = f"  [IMAGE] pos=({shape.left},{shape.top}) size=({shape.width}x{shape.height})"
            print(img_info)
        elif txt.strip():
            txt_trim = txt.strip()[:60]
            print(f'  [TEXT] "' + txt_trim + '"')
        else:
            stype = str(shape.shape_type)
            pos_info = f"  [SHAPE type={stype}] pos=({shape.left},{shape.top})"
            print(pos_info)
    print()
print("--- Verification complete ---")
