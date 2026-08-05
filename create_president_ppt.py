# -*- coding: utf-8 -*-
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
sys.path.append(r"C:\Users\FYH\AppData\Local\Programs\Python\Python313\Lib\site-packages")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import lxml.etree as ET

IMG_DIR = r"C:\Users\FYH\Documents\New project\us_imgs_hd"
OUTPUT = r"C:\Users\FYH\Documents\New project\Ding_Dongxu_President.pptx"

RED = RGBColor(205, 32, 38)
GOLD = RGBColor(212, 175, 55)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(235, 235, 240)
LTGY = RGBColor(200, 200, 210)
MDGY = RGBColor(160, 160, 170)
DKBG = RGBColor(8, 8, 35)
NAVY = RGBColor(15, 15, 50)
BLUE = RGBColor(40, 60, 130)

SW = Emu(12191365)
SH = Emu(6858000)
LI = 6
NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

def ap(sh, av):
    sp = sh._element.spPr
    sf = sp.find(NS + "solidFill")
    if sf is not None:
        sc = sf.find(NS + "srgbClr")
        if sc is not None:
            a = ET.SubElement(sc, NS + "alpha")
            a.set("val", av)

def bg(sl, ip, op="50000000"):
    sl.shapes.add_picture(ip, 0, 0, SW, SH)
    ov = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    ov.fill.solid(); ov.fill.fore_color.rgb = DKBG; ap(ov, op); ov.line.fill.background()

def tb(sl, l, t, w, h, tx, sz=18, b=False, c=WHITE, a=PP_ALIGN.LEFT):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tx
    p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = c
    p.font.name = "Microsoft YaHei"; p.alignment = a
    return bx

def ac(sl, x, y, w, c=RED, h=Emu(5000)):
    r = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = c; r.line.fill.background()

def ba(sl): ac(sl, 0, int(SH-7000), SW, RED, Emu(7000))
def ta(sl): ac(sl, 0, 0, SW, RED, Emu(7000))

def hd(sl, cn, en):
    tb(sl, Emu(600000), Emu(350000), Emu(int(SW-1200000)), Emu(480000), cn, 36, True, WHITE)
    tb(sl, Emu(600000), Emu(850000), Emu(int(SW-1200000)), Emu(280000), en, 14, False, GOLD)
    ac(sl, Emu(600000), Emu(1160000), Emu(900000), RED, Emu(5000))

def pn(sl, x, y, w, h, bg=DKBG, al="70000000"):
    p = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    p.fill.solid(); p.fill.fore_color.rgb = bg; ap(p, al)
    p.line.color.rgb = RED; p.line.width = Pt(1.5)

def cd(sl, x, y, w, h, co=RED):
    c = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = DKBG; ap(c, "70000000")
    c.line.color.rgb = co; c.line.width = Pt(1.5)
    ac(sl, x, y, w, co, Emu(5000))
    return c

def ip(n):
    p = os.path.join(IMG_DIR, n + ".jpg")
    return p if os.path.exists(p) else ""

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH

# === SLIDE 1: Title ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
p = ip("republican_rally") or ip("us_capitol")
if p: bg(s, p, "50000000")
ta(s); ba(s)
ac(s, 0, Emu(100000), SW, GOLD, Emu(2000))
ac(s, 0, int(SH-100000), SW, GOLD, Emu(2000))
tb(s, Emu(800000), Emu(800000), Emu(int(SW-1600000)), Emu(250000),
    "REPUBLICAN NATIONAL COMMITTEE  \xb7  ????????", 11, False, GOLD, PP_ALIGN.CENTER)
tb(s, Emu(800000), Emu(1200000), Emu(int(SW-1600000)), Emu(700000),
    "? ? ?", 56, True, WHITE, PP_ALIGN.CENTER)
tb(s, Emu(800000), Emu(1900000), Emu(int(SW-1600000)), Emu(400000),
    "DING DONGXU", 40, True, WHITE, PP_ALIGN.CENTER)
tb(s, Emu(800000), Emu(2500000), Emu(int(SW-1600000)), Emu(350000),
    "???????  \xb7  MAKE AMERICA GREAT AGAIN", 18, True, GOLD, PP_ALIGN.CENTER)
tb(s, Emu(800000), Emu(2950000), Emu(int(SW-1600000)), Emu(280000),
    "???????  \xb7  Republican Candidate for President", 14, False, OFF_WHITE, PP_ALIGN.CENTER)
tb(s, Emu(800000), Emu(3800000), Emu(int(SW-1600000)), Emu(200000),
    chr(9733) + "  " + chr(9733) + "  " + chr(9733) + "  " + chr(9733) + "  " + chr(9733), 16, False, GOLD, PP_ALIGN.CENTER)
tb(s, Emu(800000), Emu(4100000), Emu(int(SW-1600000)), Emu(250000),
    "???? | ???? | ??? | ?????", 13, False, LTGY, PP_ALIGN.CENTER)
tb(s, Emu(800000), Emu(4350000), Emu(int(SW-1600000)), Emu(250000),
    "Councilman | NY Governor | VP | Presidential Candidate", 11, False, MDGY, PP_ALIGN.CENTER)
cd(s, Emu(int(SW/2-2500000)), Emu(5100000), Emu(5000000), Emu(500000))
tb(s, Emu(int(SW/2-2300000)), Emu(5150000), Emu(4600000), Emu(400000),
    "??????? \xb7 MAKE AMERICA GREAT AGAIN", 16, True, GOLD, PP_ALIGN.CENTER)

# === SLIDE 2: About Me ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
p = ip("us_capitol") or ip("white_house_1")
if p: bg(s, p, "50000000")
pn(s, Emu(400000), Emu(300000), Emu(4000000), int(SH-800000), NAVY, "80000000")
tb(s, Emu(500000), Emu(500000), Emu(3800000), Emu(500000), "???", 36, True, WHITE, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(950000), Emu(3800000), Emu(300000), "ABOUT ME", 16, False, GOLD, PP_ALIGN.CENTER)
ac(s, Emu(2300000), Emu(1300000), Emu(800000), RED, Emu(5000))
circ = Emu(1700000)
cx = int(Emu(500000)+(Emu(3800000)-circ)/2)
ci = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, Emu(1500000), circ, circ)
ci.fill.solid(); ci.fill.fore_color.rgb = RGBColor(30,30,60)
ci.line.color.rgb = RED; ci.line.width = Pt(4)
tb(s, cx, Emu(2000000), circ, Emu(400000), "???", 26, True, WHITE, PP_ALIGN.CENTER)
tb(s, cx, Emu(2400000), circ, Emu(280000), "Ding Dongxu", 14, False, GOLD, PP_ALIGN.CENTER)
tb(s, cx, Emu(2680000), circ, Emu(200000), "1970 -", 12, False, LTGY, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(3650000), Emu(3800000), Emu(200000), "????  Education", 13, True, GOLD, PP_ALIGN.CENTER)
edu = ["??????  Xi'an Aviation College", "????????  SJTU Master", "??????  Harvard Ph.D."]
for i, e in enumerate(edu):
    tb(s, Emu(500000), Emu(3900000+i*250000), Emu(3800000), Emu(220000), chr(9658)+"  "+e, 11, False, OFF_WHITE)
rx = Emu(4800000); rw = Emu(int(SW-5200000))
tb(s, rx, Emu(500000), rw, Emu(400000), "????", 30, True, WHITE)
tb(s, rx, Emu(880000), rw, Emu(250000), "POLITICAL CAREER", 14, False, GOLD)
ac(s, rx, Emu(1180000), Emu(600000), RED, Emu(5000))
tlx = rx + Emu(150000)
tl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tlx, Emu(1450000), Emu(4000), Emu(4700000))
tl.fill.solid(); tl.fill.fore_color.rgb = RED; tl.line.fill.background()
career = [
    ("2000-2006", "????  Councilman", "?????????????????????"),
    ("2010-2018", "????  NY Governor", "????????????????????"),
    ("2018-2022", "???  Vice President", "???????????????????????"),
    ("2024-", "?????  Candidate", "??????????????????"),
]
for i, (per, tit, desc) in enumerate(career):
    yb = Emu(1450000+i*1200000)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, tlx-Emu(60000), yb, Emu(130000), Emu(130000))
    dot.fill.solid(); dot.fill.fore_color.rgb = GOLD; dot.line.fill.background()
    tb(s, tlx+Emu(250000), yb-Emu(50000), rw-Emu(250000), Emu(250000), per, 10, False, GOLD)
    tb(s, tlx+Emu(250000), yb+Emu(200000), rw-Emu(250000), Emu(280000), tit, 18, True, WHITE)
    tb(s, tlx+Emu(250000), yb+Emu(480000), rw-Emu(250000), Emu(400000), desc, 11, False, LTGY)
ba(s)

# === SLIDE 3: Vision ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
p = ip("statue_of_liberty") or ip("grand_canyon")
if p: bg(s, p, "50000000")
ba(s)
hd(s, "????", "MY VISION FOR AMERICA")
visions = [
    ("\u2605", "????", "Strong America", "???????????\n??????????????"),
    ("\u2605", "????", "Cooperative World", "??????????\n????????????"),
    ("\u2605", "????", "Fair Society", "??????????\n???????????"),
]
cw = Emu(3400000); ch = Emu(3500000); cgap = Emu(250000)
sx = Emu(int((SW-(cw*3+cgap*2))/2))
for i, (star, cn, en, desc) in enumerate(visions):
    cx = sx + i*(cw+cgap); cy = Emu(1600000)
    cd(s, cx, cy, cw, ch)
    tb(s, cx, cy+Emu(200000), cw, Emu(400000), star, 30, False, GOLD, PP_ALIGN.CENTER)
    tb(s, cx, cy+Emu(700000), cw, Emu(400000), cn, 28, True, WHITE, PP_ALIGN.CENTER)
    tb(s, cx, cy+Emu(1100000), cw, Emu(300000), en, 14, False, GOLD, PP_ALIGN.CENTER)
    ac(s, int(cx+cw/2-300000), cy+Emu(1500000), Emu(600000), GOLD, Emu(3000))
    tb(s, cx+Emu(300000), cy+Emu(1700000), cw-Emu(600000), Emu(1500000), desc, 14, False, LTGY, PP_ALIGN.CENTER)
tb(s, Emu(1500000), int(SH-900000), Emu(int(SW-3000000)), Emu(350000),
    chr(8220)+"????????????????????????"+chr(8221), 15, False, GOLD, PP_ALIGN.CENTER)
tb(s, Emu(1500000), int(SH-600000), Emu(int(SW-3000000)), Emu(250000),
    "\"Together, let us build a stronger, more prosperous, fairer America.\"", 12, False, LTGY, PP_ALIGN.CENTER)

# === SLIDE 4: Economy ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
p = ip("us_manufacturing") or ip("new_york_city")
if p: bg(s, p, "50000000")
ba(s)
hd(s, "??????????", "ECONOMIC POLICY: BRING MANUFACTURING BACK")
pn(s, Emu(400000), Emu(1500000), Emu(6200000), int(SH-2400000), NAVY, "75000000")
tb(s, Emu(600000), Emu(1700000), Emu(5800000), Emu(350000), chr(9755)+"  ???? / Core Policies", 18, True, GOLD)
pols = [
    ("01", "????", "Tax Incentives", "?????????????????????"),
    ("02", "????", "Jobs First", "??????????????????"),
    ("03", "????", "Energy Independence", "???????????????"),
    ("04", "????", "Fair Trade", "??????????????????"),
]
for i, (num, cn, en, desc) in enumerate(pols):
    yb = Emu(2200000+i*900000)
    ns = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(700000), yb, Emu(500000), Emu(500000))
    ns.fill.solid(); ns.fill.fore_color.rgb = RED; ns.line.fill.background()
    tb(s, Emu(700000), yb+Emu(70000), Emu(500000), Emu(380000), num, 16, True, WHITE, PP_ALIGN.CENTER)
    tb(s, Emu(1350000), yb+Emu(20000), Emu(5000000), Emu(260000), cn+"  /  "+en, 16, True, WHITE)
    tb(s, Emu(1350000), yb+Emu(280000), Emu(5000000), Emu(250000), desc, 12, False, LTGY)
r2x = Emu(6900000); r2w = Emu(int(SW-7400000))
pn(s, r2x, Emu(1500000), r2w, int(SH-2400000), DKBG, "80000000")
tb(s, r2x+Emu(200000), Emu(1700000), r2w-Emu(400000), Emu(250000), "?? / Goals", 16, True, GOLD, PP_ALIGN.CENTER)
ac(s, int(r2x+r2w/2-300000), Emu(2000000), Emu(600000), RED, Emu(3000))
stats = [("1000?+", "????\nNew Jobs"), ("20%", "????\nCorporate Tax"), ("5??$", "GDP??\nGDP Growth")]
for i, (num, label) in enumerate(stats):
    ys = Emu(2300000+i*1200000)
    tb(s, r2x+Emu(200000), ys, r2w-Emu(400000), Emu(350000), num, 36, True, GOLD, PP_ALIGN.CENTER)
    tb(s, r2x+Emu(200000), ys+Emu(380000), r2w-Emu(400000), Emu(250000), label, 13, True, WHITE, PP_ALIGN.CENTER)

# === SLIDE 5: Immigration ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
p = ip("statue_of_liberty") or ip("liberty_bell")
if p: bg(s, p, "50000000")
ba(s)
hd(s, "??????????", "IMMIGRATION POLICY: ORDER AND SECURITY")
pn(s, Emu(400000), Emu(1500000), int(SW-800000), int(SH-2400000), NAVY, "75000000")
col_w = Emu(int((SW-1800000)/2))
tb(s, Emu(600000), Emu(1700000), col_w, Emu(300000), chr(9940)+"  ???? / Restrict Immigration", 18, True, RED)
left = ["???????? / Ban illegal immigration", "?????????????", "??????????", "??????", "??????"]
for i, it in enumerate(left):
    tb(s, Emu(750000), Emu(2100000+i*320000), col_w+Emu(500000), Emu(280000), chr(10006)+"  "+it, 12, False, OFF_WHITE)
rx2 = Emu(6600000)
tb(s, rx2, Emu(1700000), col_w, Emu(300000), chr(10004)+"  ???? / Orderly Management", 18, True, GOLD)
right = ["???????? / Optimize legal channels", "??????? / Attract skilled talent",
    "???????????", "?????????", "????????"]
for i, it in enumerate(right):
    tb(s, rx2+Emu(100000), Emu(2100000+i*320000), col_w+Emu(500000), Emu(280000), chr(10004)+"  "+it, 12, False, OFF_WHITE)
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(SW/2-2000), Emu(2100000), Emu(4000), Emu(1600000))
div.fill.solid(); div.fill.fore_color.rgb = RED; div.line.fill.background()
tb(s, Emu(800000), int(SH-800000), Emu(int(SW-1600000)), Emu(300000),
    chr(8220)+"???????????????"+chr(8212)+"????"+chr(8221), 14, False, GOLD, PP_ALIGN.CENTER)
tb(s, Emu(800000), int(SH-550000), Emu(int(SW-1600000)), Emu(250000),
    "\"Legal immigrants welcome, illegal immigrants not tolerated.\"", 11, False, LTGY, PP_ALIGN.CENTER)

# === SLIDE 6: Foreign Policy ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
p = ip("us_china_cooperation") or ip("times_square")
if p: bg(s, p, "50000000")
ba(s)
hd(s, "??????????", "FOREIGN POLICY: PRAGMATISM AND COOPERATION")
pn(s, Emu(300000), Emu(1500000), int(SW-600000), int(SH-2400000), NAVY, "75000000")
tb(s, Emu(500000), Emu(1700000), int(SW-1200000), Emu(300000),
    chr(9755)+"  ???????? / Reshaping America's Foreign Relations", 18, True, GOLD)
pillars = [
    ("\U0001f30d", "????", "China-US Cooperation",
     "????????????????\n????????????????",
     "Bipolar world. Cooperate with China."),
    ("\U0001f30e", "??????", "Realign Alliances",
     "???????????????\n????????????????",
     "Reduce EU support, reassess SK/Japan aid."),
    ("\U0001f30f", "????", "America First",
     "????????????\n??????????????",
     "National interest first. Focus inward."),
]
pw = Emu(3600000); ph = Emu(2800000); pgap = Emu(250000)
psx = Emu(int(450000))
for i, (icon, cn, en, cnd, end) in enumerate(pillars):
    px = psx + i*(pw+pgap)
    cd(s, px, Emu(2200000), pw, ph, RED if i==0 else (GOLD if i==1 else BLUE))
    tb(s, px, Emu(2350000), pw, Emu(400000), icon, 28, False, WHITE, PP_ALIGN.CENTER)
    tb(s, px, Emu(2750000), pw, Emu(300000), cn, 20, True, WHITE, PP_ALIGN.CENTER)
    tb(s, px, Emu(3050000), pw, Emu(250000), en, 13, False, GOLD, PP_ALIGN.CENTER)
    ac(s, int(px+pw/2-250000), Emu(3350000), Emu(500000), GOLD, Emu(3000))
    tb(s, px+Emu(200000), Emu(3450000), pw-Emu(400000), Emu(1300000), cnd, 12, False, OFF_WHITE)
    tb(s, px+Emu(200000), Emu(4600000), pw-Emu(400000), Emu(350000), end, 10, False, MDGY)

# === SLIDE 7: Infrastructure & Healthcare ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
p = ip("us_infrastructure") or ip("american_healthcare")
if p: bg(s, p, "50000000")
ba(s)
hd(s, "???????", "INFRASTRUCTURE & UNIVERSAL HEALTHCARE")
pn(s, Emu(400000), Emu(1500000), int(SW-800000), Emu(1600000), NAVY, "75000000")
tb(s, Emu(600000), Emu(1650000), Emu(5000000), Emu(300000),
    "\U0001f3d7\ufe0f  ?????? / Infrastructure Rebuild", 18, True, GOLD)
for i, it in enumerate([
    chr(9658)+"  ??2????????????????? /  for roads, bridges, airports",
    chr(9658)+"  ??????????? / Build modern high-speed rail network",
    chr(9658)+"  ??????????? / Upgrade power grid and broadband",
    chr(9658)+"  ???????????? / Promote green energy infrastructure",
]):
    tb(s, Emu(700000), Emu(2050000+i*280000), Emu(5500000), Emu(250000), it, 12, False, OFF_WHITE)
rhx = Emu(6500000)
tb(s, rhx, Emu(1650000), Emu(5000000), Emu(300000),
    "\U0001f3e5  ???? / Universal Healthcare", 18, True, RED)
for i, it in enumerate([
    chr(9658)+"  ???????? / Push for universal health coverage",
    chr(9658)+"  ??????? / Lower prescription drug prices",
    chr(9658)+"  ?????????????? / Tax big corps to fund healthcare",
    chr(9658)+"  ?????????? / Strengthen medical infrastructure",
]):
    tb(s, rhx+Emu(100000), Emu(2050000+i*280000), Emu(5500000), Emu(250000), it, 12, False, OFF_WHITE)
div2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(SW/2-2000), Emu(2000000), Emu(4000), Emu(1200000))
div2.fill.solid(); div2.fill.fore_color.rgb = GOLD; div2.line.fill.background()
for i, (num, lab) in enumerate([("", "????\nInfrastructure"), ("500?", "????\nNew Jobs"), ("??", "????\nHealthcare")]):
    sx = Emu(1200000+i*3300000)
    tb(s, sx, Emu(3400000), Emu(2800000), Emu(350000), num, 32, True, GOLD, PP_ALIGN.CENTER)
    tb(s, sx, Emu(3750000), Emu(2800000), Emu(250000), lab, 13, False, LTGY, PP_ALIGN.CENTER)

# === SLIDE 8: Education ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
p = ip("us_education") or ip("times_square")
if p: bg(s, p, "50000000")
ba(s)
hd(s, "??????????????", "EDUCATION REFORM: COMPULSORY EDUCATION & EXAMS")
pn(s, Emu(400000), Emu(1500000), int(SW-800000), int(SH-2400000), NAVY, "75000000")
tb(s, Emu(600000), Emu(1700000), int(SW-1200000), Emu(300000),
    chr(9755)+"  ?????? / Key Reform Measures", 18, True, GOLD)
reforms = [
    ("\U0001f393", "????????", "9-Year Compulsory Education",
     "????????????????????????????"),
    ("\U0001f4da", "??????", "Exam Selection System",
     "????????????????????????????"),
    ("\U0001f3eb", "??????", "Quality Improvement",
     "??????????????????????????"),
    ("\U0001f4bc", "??????", "Vocational Education",
     "?????????????????????"),
]
for i, (icon, cn, en, desc) in enumerate(reforms):
    col = i%2; row = i//2
    ix = Emu(600000)+col*Emu(5600000); iy = Emu(2150000)+row*Emu(1450000)
    tb(s, ix, iy, Emu(500000), Emu(400000), icon, 24, False, WHITE)
    tb(s, ix+Emu(500000), iy, Emu(5000000), Emu(300000), cn+"  /  "+en, 16, True, WHITE)
    tb(s, ix+Emu(500000), iy+Emu(350000), Emu(5000000), Emu(400000), desc, 11, False, OFF_WHITE)

# === SLIDE 9: Why Vote for Me ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
p = ip("mount_rushmore") or ip("republican_rally")
if p: bg(s, p, "50000000")
ba(s)
hd(s, "?????", "WHY VOTE FOR DING DONGXU")
pn(s, Emu(400000), Emu(1500000), int(SW-800000), int(SH-2400000), NAVY, "75000000")
reasons = [
    ("\U0001f3db\ufe0f", "????", "Proven Experience", "?????????\n30????????"),
    ("\U0001f1fa\U0001f1f8", "????", "Deeply Patriotic", "????????\n????????"),
    ("\U0001f4aa", "????", "Results-Oriented", "????????????\n?????????"),
    ("\U0001f30a", "????", "Visionary", "????????\n????????"),
]
rw2 = Emu(2700000); rh2 = Emu(2100000); rgap = Emu(250000)
rsx = Emu(550000)
for i, (icon, cn, en, desc) in enumerate(reasons):
    col = i%2; row = i//2
    rx = rsx+col*(rw2+rgap+Emu(100000)); ry = Emu(1750000)+row*(rh2+Emu(150000))
    cd(s, rx, ry, rw2, rh2)
    tb(s, rx, ry+Emu(150000), rw2, Emu(350000), icon, 24, False, WHITE, PP_ALIGN.CENTER)
    tb(s, rx, ry+Emu(500000), rw2, Emu(300000), cn, 20, True, WHITE, PP_ALIGN.CENTER)
    tb(s, rx, ry+Emu(800000), rw2, Emu(250000), en, 12, False, GOLD, PP_ALIGN.CENTER)
    ac(s, int(rx+rw2/2-250000), ry+Emu(1100000), Emu(500000), GOLD, Emu(3000))
    tb(s, rx+Emu(150000), ry+Emu(1200000), rw2-Emu(300000), Emu(700000), desc, 12, False, LTGY, PP_ALIGN.CENTER)

# === SLIDE 10: Call to Action ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
p = ip("american_flag_eagle") or ip("white_house_1")
if p: bg(s, p, "45000000")
ta(s); ba(s)
ac(s, 0, Emu(100000), SW, GOLD, Emu(2000))
ac(s, 0, int(SH-100000), SW, GOLD, Emu(2000))
tb(s, Emu(500000), Emu(1200000), Emu(int(SW-1000000)), Emu(600000),
    "????????", 48, True, WHITE, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(1750000), Emu(int(SW-1000000)), Emu(350000),
    "MAKE AMERICA GREAT AGAIN!", 36, True, GOLD, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(2300000), Emu(int(SW-1000000)), Emu(300000),
    chr(9733)+" \xb7 "+chr(9733)+" \xb7 "+chr(9733)+" \xb7 "+chr(9733)+" \xb7 "+chr(9733), 18, False, RED, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(2700000), Emu(int(SW-1000000)), Emu(350000),
    "????? \u2014 ???", 28, True, WHITE, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(3050000), Emu(int(SW-1000000)), Emu(300000),
    "VOTE FOR ME \u2014 DING DONGXU FOR PRESIDENT", 18, True, GOLD, PP_ALIGN.CENTER)
cd(s, Emu(int(SW/2-2500000)), Emu(3600000), Emu(5000000), Emu(650000), GOLD)
tb(s, Emu(int(SW/2-2300000)), Emu(3680000), Emu(4600000), Emu(260000),
    chr(8220)+"???? \xb7 ???? \xb7 ???????"+chr(8221), 16, True, WHITE, PP_ALIGN.CENTER)
tb(s, Emu(int(SW/2-2300000)), Emu(3950000), Emu(4600000), Emu(250000),
    "\"Serve with heart. Lead with action. Make America Great Again.\"", 12, False, GOLD, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(4650000), Emu(int(SW-1000000)), Emu(250000),
    "2024??????? \xb7 ?????? \xb7 ???", 13, False, LTGY, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(4880000), Emu(int(SW-1000000)), Emu(250000),
    "2024 U.S. Presidential Election \xb7 Republican Candidate \xb7 Ding Dongxu", 11, False, MDGY, PP_ALIGN.CENTER)

# === SAVE ===
prs.save(OUTPUT)
print("PPT generated successfully!")
print("Output: " + OUTPUT)
print("Total slides: " + str(len(prs.slides)))
