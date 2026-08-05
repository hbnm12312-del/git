import json, sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import lxml.etree as ET
import json, sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import lxml.etree as ET
# Load dictionary
_ = json.load(open(r'C:\Users\FYH\Documents\New project\texts.json', 'r', encoding='utf-8'))
# Constants
IMG_DIR = r'C:\Users\FYH\Documents\New project\us_imgs_hd'
OUTPUT = r'C:\Users\FYH\Documents\New project\Ding_Dongxu_President.pptx'
RED = RGBColor(205, 32, 38); GOLD = RGBColor(212, 175, 55)
WHITE = RGBColor(255, 255, 255); OFF_WHITE = RGBColor(235, 235, 240)
LTGY = RGBColor(200, 200, 210); MDGY = RGBColor(160, 160, 170)
DKBG = RGBColor(8, 8, 35); NAVY = RGBColor(15, 15, 50); BLUE = RGBColor(40, 60, 130)
SW = Emu(12191365); SH = Emu(6858000); LI = 6
NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
def ap(sh, av):
    sp = sh._element.spPr; sf = sp.find(NS + 'solidFill')
    if sf is not None:
        sc = sf.find(NS + 'srgbClr')
        if sc is not None:
            a = ET.SubElement(sc, NS + 'alpha'); a.set('val', av)
def bg(sl, ip, op='12000000'):
    sl.shapes.add_picture(ip, 0, 0, SW, SH)
    ov = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    ov.fill.solid(); ov.fill.fore_color.rgb = DKBG; ap(ov, op); ov.line.fill.background()
def tb(sl, l, t, w, h, tx, sz=18, b=False, c=WHITE, a=PP_ALIGN.LEFT):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tx
    p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = c
    p.font.name = 'Microsoft YaHei'; p.alignment = a; return bx
def ac(sl, x, y, w, c=RED, h=Emu(5000)):
    r = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = c; r.line.fill.background()
def ba(sl): ac(sl, 0, int(SH-7000), SW, RED, Emu(7000))
def ta(sl): ac(sl, 0, 0, SW, RED, Emu(7000))
def hd(sl, cn, en):
    tb(sl, Emu(600000), Emu(350000), Emu(int(SW-1200000)), Emu(480000), cn, 36, True, WHITE)
    tb(sl, Emu(600000), Emu(850000), Emu(int(SW-1200000)), Emu(280000), en, 14, False, GOLD)
    ac(sl, Emu(600000), Emu(1160000), Emu(900000), RED, Emu(5000))
def pn(sl, x, y, w, h, bg=DKBG, al='70000000'):
    p = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    p.fill.solid(); p.fill.fore_color.rgb = bg; ap(p, al)
    p.line.color.rgb = RED; p.line.width = Pt(1.5)
def cd(sl, x, y, w, h, co=RED):
    c = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = DKBG; ap(c, '70000000')
    c.line.color.rgb = co; c.line.width = Pt(1.5)
    ac(sl, x, y, w, co, Emu(5000)); return c
def ip(n):
    p = os.path.join(IMG_DIR, n + '.jpg')
    return p if os.path.exists(p) else ''
def sb(s, key):
    p = ip(key)
    if p: bg(s, p, '50000000')
# === SLIDE 1: Title ===
def sb(s, key, op='20000000'):
    p = ip(key)
    if p: bg(s, p, op)
# === SLIDE 1: Title ===
def sb(s, key, op='20000000'):
    p = ip(key)
    if p:
        bg(s, p, op)
        return True
    return False
def bo(s, key, op='20000000'):
    p = ip(key)
    if p:
        bg(s, p, op)
        return True
    if key.find(',') > 0:
        for k in key.split(','):
            p = ip(k.strip())
            if p:
                bg(s, p, op)
                return True
    return False
# Create presentation
prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
print('Script ready - generating...')
print('Dictionary loaded:', len(_), 'keys')
print('Dictionary loaded:', len(_), 'keys')
# === SLIDE 1: Title ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
bo(s, 'republican_rally,us_capitol', '12000000')
ta(s); ba(s)
ac(s, 0, Emu(100000), SW, GOLD, Emu(2000))
ac(s, 0, int(SH-100000), SW, GOLD, Emu(2000))
tb(s, Emu(800000), Emu(800000), Emu(int(SW-1600000)), Emu(250000), _['s1_gop'], 11, False, GOLD, PP_ALIGN.CENTER)
tb(s, Emu(800000), Emu(1200000), Emu(int(SW-1600000)), Emu(700000), _['s1_name_cn'], 56, True, WHITE, PP_ALIGN.CENTER)
tb(s, Emu(800000), Emu(1900000), Emu(int(SW-1600000)), Emu(400000), _['s1_name_en'], 40, True, WHITE, PP_ALIGN.CENTER)
tb(s, Emu(800000), Emu(2500000), Emu(int(SW-1600000)), Emu(350000), _['s1_slogan'], 18, True, GOLD, PP_ALIGN.CENTER)
tb(s, Emu(800000), Emu(2950000), Emu(int(SW-1600000)), Emu(280000), _['s1_sub'], 14, False, OFF_WHITE, PP_ALIGN.CENTER)
st = '\u2605 \u00a0 \u2605 \u00a0 \u2605 \u00a0 \u2605 \u00a0 \u2605'
tb(s, Emu(800000), Emu(3800000), Emu(int(SW-1600000)), Emu(200000), st, 16, False, GOLD, PP_ALIGN.CENTER)
tb(s, Emu(800000), Emu(4100000), Emu(int(SW-1600000)), Emu(250000), _['s1_career_cn'], 13, False, LTGY, PP_ALIGN.CENTER)
tb(s, Emu(800000), Emu(4350000), Emu(int(SW-1600000)), Emu(250000), _['s1_career_en'], 11, False, MDGY, PP_ALIGN.CENTER)
cd(s, Emu(int(SW/2-2500000)), Emu(5100000), Emu(5000000), Emu(500000))
tb(s, Emu(int(SW/2-2300000)), Emu(5150000), Emu(4600000), Emu(400000), _['s1_bottom'], 16, True, GOLD, PP_ALIGN.CENTER)
print('Slide 1 done')
# === SLIDE 2: About Me ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
bo(s, 'us_capitol,white_house', '05000000')
pn(s, Emu(400000), Emu(300000), Emu(4000000), int(SH-800000), NAVY, '80000000')
tb(s, Emu(500000), Emu(500000), Emu(3800000), Emu(500000), _['s2_title_cn'], 36, True, WHITE, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(950000), Emu(3800000), Emu(300000), _['s2_title_en'], 16, False, GOLD, PP_ALIGN.CENTER)
ac(s, Emu(2300000), Emu(1300000), Emu(800000), RED, Emu(5000))
circ = Emu(1700000); cx = int(Emu(500000)+(Emu(3800000)-circ)/2)
ci = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, Emu(1500000), circ, circ)
ci.fill.solid(); ci.fill.fore_color.rgb = RGBColor(30,30,60); ci.line.color.rgb = RED; ci.line.width = Pt(4)
tb(s, cx, Emu(2000000), circ, Emu(400000), _['s2_name'], 26, True, WHITE, PP_ALIGN.CENTER)
tb(s, cx, Emu(2400000), circ, Emu(280000), 'Ding Dongxu', 14, False, GOLD, PP_ALIGN.CENTER)
tb(s, cx, Emu(2680000), circ, Emu(200000), '1970 -', 12, False, LTGY, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(3650000), Emu(3800000), Emu(200000), _['s2_edu_title'], 13, True, GOLD, PP_ALIGN.CENTER)
edu = [_['s2_edu_1'], _['s2_edu_2'], _['s2_edu_3']]
for i, e in enumerate(edu):
    tb(s, Emu(500000), Emu(3900000+i*250000), Emu(3800000), Emu(220000), '\u25b8  '+e, 11, False, OFF_WHITE)
rx = Emu(4800000); rw = Emu(int(SW-5200000))
tb(s, rx, Emu(500000), rw, Emu(400000), _['s2_career_title_cn'], 30, True, WHITE)
tb(s, rx, Emu(880000), rw, Emu(250000), _['s2_career_title_en'], 14, False, GOLD)
ac(s, rx, Emu(1180000), Emu(600000), RED, Emu(5000))
tlx = rx + Emu(150000)
tl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tlx, Emu(1450000), Emu(4000), Emu(4700000))
tl.fill.solid(); tl.fill.fore_color.rgb = RED; tl.line.fill.background()
career = [('s2_c1', _['s2_c1_per'], _['s2_c1_tit'], _['s2_c1_desc']),
          ('s2_c2', _['s2_c2_per'], _['s2_c2_tit'], _['s2_c2_desc']),
          ('s2_c3', _['s2_c3_per'], _['s2_c3_tit'], _['s2_c3_desc']),
          ('s2_c4', _['s2_c4_per'], _['s2_c4_tit'], _['s2_c4_desc'])]
for i, (k, per, tit, desc) in enumerate(career):
    yb = Emu(1450000+i*1200000)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, tlx-Emu(60000), yb, Emu(130000), Emu(130000))
    dot.fill.solid(); dot.fill.fore_color.rgb = GOLD; dot.line.fill.background()
    tb(s, tlx+Emu(250000), yb-Emu(50000), rw-Emu(250000), Emu(250000), per, 10, False, GOLD)
    tb(s, tlx+Emu(250000), yb+Emu(200000), rw-Emu(250000), Emu(280000), tit, 18, True, WHITE)
    tb(s, tlx+Emu(250000), yb+Emu(480000), rw-Emu(250000), Emu(400000), desc, 11, False, LTGY)
ba(s)
print('Slide 2 done')
print('Slide 2 done')
# === SLIDE 3: Vision ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
bo(s, 'statue_of_liberty,grand_canyon', '12000000')
ba(s)
hd(s, _['s3_title_cn'], _['s3_title_en'])
visions = [('\u2605', 'v1', _['s3_v1_cn'], _['s3_v1_en'], _['s3_v1_desc']),
           ('\u2605', 'v2', _['s3_v2_cn'], _['s3_v2_en'], _['s3_v2_desc']),
           ('\u2605', 'v3', _['s3_v3_cn'], _['s3_v3_en'], _['s3_v3_desc'])]
cw = Emu(3400000); ch = Emu(3500000); cgap = Emu(250000)
sx = Emu(int((SW-(cw*3+cgap*2))/2))
for i, (star, k, cn, en, desc) in enumerate(visions):
    cx = sx + i*(cw+cgap); cy = Emu(1600000)
    cd(s, cx, cy, cw, ch)
    tb(s, cx, cy+Emu(200000), cw, Emu(400000), star, 30, False, GOLD, PP_ALIGN.CENTER)
    tb(s, cx, cy+Emu(700000), cw, Emu(400000), cn, 28, True, WHITE, PP_ALIGN.CENTER)
    tb(s, cx, cy+Emu(1100000), cw, Emu(300000), en, 14, False, GOLD, PP_ALIGN.CENTER)
    ac(s, int(cx+cw/2-300000), cy+Emu(1500000), Emu(600000), GOLD, Emu(3000))
    tb(s, cx+Emu(300000), cy+Emu(1700000), cw-Emu(600000), Emu(1500000), desc, 14, False, LTGY, PP_ALIGN.CENTER)
tb(s, Emu(1500000), int(SH-900000), Emu(int(SW-3000000)), Emu(350000), '\u201c'+_['s3_quote']+'\u201d', 15, False, GOLD, PP_ALIGN.CENTER)
tb(s, Emu(1500000), int(SH-600000), Emu(int(SW-3000000)), Emu(250000), '\u201c'+_['s3_quote_en']+'\u201d', 12, False, LTGY, PP_ALIGN.CENTER)
print('Slide 3 done')
# === SLIDE 4: Economy ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
bo(s, 'us_manufacturing,new_york_city', '12000000')
ba(s)
hd(s, _['s4_title_cn'], _['s4_title_en'])
pn(s, Emu(400000), Emu(1500000), Emu(6200000), int(SH-2400000), NAVY, '75000000')
tb(s, Emu(600000), Emu(1700000), Emu(5800000), Emu(350000), '\u261b  '+_['s4_sub'], 18, True, GOLD)
pols = [('01', 'p1', _['s4_p1_cn'], _['s4_p1_en'], _['s4_p1_desc']),
         ('02', 'p2', _['s4_p2_cn'], _['s4_p2_en'], _['s4_p2_desc']),
         ('03', 'p3', _['s4_p3_cn'], _['s4_p3_en'], _['s4_p3_desc']),
         ('04', 'p4', _['s4_p4_cn'], _['s4_p4_en'], _['s4_p4_desc'])]
for i, (num, k, cn, en, desc) in enumerate(pols):
    yb = Emu(2200000+i*900000)
    ns = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(700000), yb, Emu(500000), Emu(500000))
    ns.fill.solid(); ns.fill.fore_color.rgb = RED; ns.line.fill.background()
    tb(s, Emu(700000), yb+Emu(70000), Emu(500000), Emu(380000), num, 16, True, WHITE, PP_ALIGN.CENTER)
    tb(s, Emu(1350000), yb+Emu(20000), Emu(5000000), Emu(260000), cn+'  /  '+en, 16, True, WHITE)
    tb(s, Emu(1350000), yb+Emu(280000), Emu(5000000), Emu(250000), desc, 12, False, LTGY)
r2x = Emu(6900000); r2w = Emu(int(SW-7400000))
pn(s, r2x, Emu(1500000), r2w, int(SH-2400000), DKBG, '80000000')
tb(s, r2x+Emu(200000), Emu(1700000), r2w-Emu(400000), Emu(250000), _['s4_goal'], 16, True, GOLD, PP_ALIGN.CENTER)
ac(s, int(r2x+r2w/2-300000), Emu(2000000), Emu(600000), RED, Emu(3000))
stats = [('s1', _['s4_s1_num'], _['s4_s1_lab']),
         ('s2', _['s4_s2_num'], _['s4_s2_lab']),
         ('s3', _['s4_s3_num'], _['s4_s3_lab'])]
for i, (k, num, lab) in enumerate(stats):
    ys = Emu(2300000+i*1200000)
    tb(s, r2x+Emu(200000), ys, r2w-Emu(400000), Emu(350000), num, 36, True, GOLD, PP_ALIGN.CENTER)
    tb(s, r2x+Emu(200000), ys+Emu(380000), r2w-Emu(400000), Emu(250000), lab, 13, True, WHITE, PP_ALIGN.CENTER)
print('Slide 4 done')
# === SLIDE 5: Immigration ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
bo(s, 'statue_of_liberty,liberty_bell', '12000000')
ba(s)
hd(s, _['s5_title_cn'], _['s5_title_en'])
pn(s, Emu(400000), Emu(1500000), int(SW-800000), int(SH-2400000), NAVY, '75000000')
col_w = Emu(int((SW-1800000)/2))
tb(s, Emu(600000), Emu(1700000), col_w, Emu(300000), '\u26d4  '+_['s5_left_title'], 18, True, RED)
left = [_['s5_l1'], _['s5_l2'], _['s5_l3'], _['s5_l4'], _['s5_l5']]
for i, it in enumerate(left):
    tb(s, Emu(750000), Emu(2100000+i*320000), col_w+Emu(500000), Emu(280000), '\u2716  '+it, 12, False, OFF_WHITE)
rx2 = Emu(6600000)
tb(s, rx2, Emu(1700000), col_w, Emu(300000), '\u2714\ufe0f  '+_['s5_right_title'], 18, True, GOLD)
right = [_['s5_r1'], _['s5_r2'], _['s5_r3'], _['s5_r4'], _['s5_r5']]
for i, it in enumerate(right):
    tb(s, rx2+Emu(100000), Emu(2100000+i*320000), col_w+Emu(500000), Emu(280000), '\u2714  '+it, 12, False, OFF_WHITE)
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(SW/2-2000), Emu(2100000), Emu(4000), Emu(1600000))
div.fill.solid(); div.fill.fore_color.rgb = RED; div.line.fill.background()
tb(s, Emu(800000), int(SH-800000), Emu(int(SW-1600000)), Emu(300000), '\u201c'+_['s5_quote']+'\u201d', 14, False, GOLD, PP_ALIGN.CENTER)
tb(s, Emu(800000), int(SH-550000), Emu(int(SW-1600000)), Emu(250000), '\u201c'+_['s5_quote_en']+'\u201d', 11, False, LTGY, PP_ALIGN.CENTER)
print('Slide 5 done')
print('Slide 5 done')
# === SLIDE 6: Foreign Policy ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
bo(s, 'us_china_cooperation,times_square', '05000000')
ba(s)
hd(s, _['s6_title_cn'], _['s6_title_en'])
pn(s, Emu(300000), Emu(1500000), int(SW-600000), int(SH-2400000), NAVY, '75000000')
tb(s, Emu(500000), Emu(1700000), int(SW-1200000), Emu(300000), '\u261b  '+_['s6_sub'], 18, True, GOLD)
pillars = [('\U0001f30d', 'p1', _['s6_p1_cn'], _['s6_p1_en'], _['s6_p1_desc'], _['s6_p1_en_desc']),
           ('\U0001f30e', 'p2', _['s6_p2_cn'], _['s6_p2_en'], _['s6_p2_desc'], _['s6_p2_en_desc']),
           ('\U0001f30f', 'p3', _['s6_p3_cn'], _['s6_p3_en'], _['s6_p3_desc'], _['s6_p3_en_desc'])]
pw = Emu(3600000); ph = Emu(2800000); pgap = Emu(250000)
psx = Emu(int(450000))
for i, (icon, k, cn, en, cnd, end) in enumerate(pillars):
    px = psx + i*(pw+pgap)
    co = RED if i==0 else (GOLD if i==1 else BLUE)
    cd(s, px, Emu(2200000), pw, ph, co)
    tb(s, px, Emu(2350000), pw, Emu(400000), icon, 28, False, WHITE, PP_ALIGN.CENTER)
    tb(s, px, Emu(2750000), pw, Emu(300000), cn, 20, True, WHITE, PP_ALIGN.CENTER)
    tb(s, px, Emu(3050000), pw, Emu(250000), en, 13, False, GOLD, PP_ALIGN.CENTER)
    ac(s, int(px+pw/2-250000), Emu(3350000), Emu(500000), GOLD, Emu(3000))
    tb(s, px+Emu(200000), Emu(3450000), pw-Emu(400000), Emu(1300000), cnd, 12, False, OFF_WHITE)
    tb(s, px+Emu(200000), Emu(4600000), pw-Emu(400000), Emu(350000), end, 10, False, MDGY)
print('Slide 6 done')
# === SLIDE 7: Infrastructure & Healthcare ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
bo(s, 'us_infrastructure,american_healthcare', '12000000')
ba(s)
hd(s, _['s7_title_cn'], _['s7_title_en'])
pn(s, Emu(400000), Emu(1500000), int(SW-800000), Emu(1600000), NAVY, '75000000')
tb(s, Emu(600000), Emu(1650000), Emu(5000000), Emu(300000), '\U0001f3d7\ufe0f  '+_['s7_infra_title'], 18, True, GOLD)
infra = [_['s7_infra_1'], _['s7_infra_2'], _['s7_infra_3'], _['s7_infra_4']]
for i, it in enumerate(infra):
    tb(s, Emu(700000), Emu(2050000+i*280000), Emu(5500000), Emu(250000), '\u25b8  '+it, 12, False, OFF_WHITE)
rhx = Emu(6500000)
tb(s, rhx, Emu(1650000), Emu(5000000), Emu(300000), '\U0001f3e5  '+_['s7_health_title'], 18, True, RED)
health = [_['s7_health_1'], _['s7_health_2'], _['s7_health_3'], _['s7_health_4']]
for i, it in enumerate(health):
    tb(s, rhx+Emu(100000), Emu(2050000+i*280000), Emu(5500000), Emu(250000), '\u25b8  '+it, 12, False, OFF_WHITE)
div2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(SW/2-2000), Emu(2000000), Emu(4000), Emu(1200000))
div2.fill.solid(); div2.fill.fore_color.rgb = GOLD; div2.line.fill.background()
for i, (k, num, lab) in enumerate([('s1', _['s7_s1_num'], _['s7_s1_lab']),
    ('s2', _['s7_s2_num'], _['s7_s2_lab']), ('s3', _['s7_s3_num'], _['s7_s3_lab'])]):
    sx = Emu(1200000+i*3300000)
    tb(s, sx, Emu(3400000), Emu(2800000), Emu(350000), num, 32, True, GOLD, PP_ALIGN.CENTER)
    tb(s, sx, Emu(3750000), Emu(2800000), Emu(250000), lab, 13, False, LTGY, PP_ALIGN.CENTER)
print('Slide 7 done')
# === SLIDE 8: Education ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
bo(s, 'us_education,times_square', '12000000')
ba(s)
hd(s, _['s8_title_cn'], _['s8_title_en'])
pn(s, Emu(400000), Emu(1500000), int(SW-800000), int(SH-2400000), NAVY, '75000000')
tb(s, Emu(600000), Emu(1700000), int(SW-1200000), Emu(300000), '\u261b  '+_['s8_sub'], 18, True, GOLD)
reforms = [('\U0001f393', 'r1', _['s8_r1_cn'], _['s8_r1_en'], _['s8_r1_desc']),
           ('\U0001f4da', 'r2', _['s8_r2_cn'], _['s8_r2_en'], _['s8_r2_desc']),
           ('\U0001f3eb', 'r3', _['s8_r3_cn'], _['s8_r3_en'], _['s8_r3_desc']),
           ('\U0001f4bc', 'r4', _['s8_r4_cn'], _['s8_r4_en'], _['s8_r4_desc'])]
for i, (icon, k, cn, en, desc) in enumerate(reforms):
    col = i%2; row = i//2
    ix = Emu(600000)+col*Emu(5600000); iy = Emu(2150000)+row*Emu(1450000)
    tb(s, ix, iy, Emu(500000), Emu(400000), icon, 24, False, WHITE)
    tb(s, ix+Emu(500000), iy, Emu(5000000), Emu(300000), cn+'  /  '+en, 16, True, WHITE)
    tb(s, ix+Emu(500000), iy+Emu(350000), Emu(5000000), Emu(400000), desc, 11, False, OFF_WHITE)
print('Slide 8 done')
print('Slide 8 done')
# === SLIDE 9: Why Vote ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
bo(s, 'mount_rushmore,republican_rally', '08000000')
ba(s)
hd(s, _['s9_title_cn'], _['s9_title_en'])
pn(s, Emu(400000), Emu(1500000), int(SW-800000), int(SH-2400000), NAVY, '75000000')
reasons = [('\U0001f3db\ufe0f', 'r1', _['s9_r1_cn'], _['s9_r1_en'], _['s9_r1_desc']),
           ('\U0001f1fa\U0001f1f8', 'r2', _['s9_r2_cn'], _['s9_r2_en'], _['s9_r2_desc']),
           ('\U0001f4aa', 'r3', _['s9_r3_cn'], _['s9_r3_en'], _['s9_r3_desc']),
           ('\U0001f30a', 'r4', _['s9_r4_cn'], _['s9_r4_en'], _['s9_r4_desc'])]
rw = Emu(2700000); rh = Emu(2100000); rgap = Emu(250000); rsx = Emu(550000)
for i, (icon, k, cn, en, desc) in enumerate(reasons):
    col = i%2; row = i//2
    rx = rsx+col*(rw+rgap+Emu(100000)); ry = Emu(1750000)+row*(rh+Emu(150000))
    cd(s, rx, ry, rw, rh)
    tb(s, rx, ry+Emu(150000), rw, Emu(350000), icon, 24, False, WHITE, PP_ALIGN.CENTER)
    tb(s, rx, ry+Emu(500000), rw, Emu(300000), cn, 20, True, WHITE, PP_ALIGN.CENTER)
    tb(s, rx, ry+Emu(800000), rw, Emu(250000), en, 12, False, GOLD, PP_ALIGN.CENTER)
    ac(s, int(rx+rw/2-250000), ry+Emu(1100000), Emu(500000), GOLD, Emu(3000))
    tb(s, rx+Emu(150000), ry+Emu(1200000), rw-Emu(300000), Emu(700000), desc, 12, False, LTGY, PP_ALIGN.CENTER)
print('Slide 9 done')
# === SLIDE 10: Call to Action ===
s = prs.slides.add_slide(prs.slide_layouts[LI])
bo(s, 'american_flag_eagle,white_house', '12000000')
ta(s); ba(s)
ac(s, 0, Emu(100000), SW, GOLD, Emu(2000))
ac(s, 0, int(SH-100000), SW, GOLD, Emu(2000))
tb(s, Emu(500000), Emu(1200000), Emu(int(SW-1000000)), Emu(600000), _['s10_slogan'], 48, True, WHITE, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(1750000), Emu(int(SW-1000000)), Emu(350000), _['s10_slogan_en'], 36, True, GOLD, PP_ALIGN.CENTER)
st = '\u2605 \u00b7 \u2605 \u00b7 \u2605 \u00b7 \u2605 \u00b7 \u2605'
tb(s, Emu(500000), Emu(2300000), Emu(int(SW-1000000)), Emu(300000), st, 18, False, RED, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(2700000), Emu(int(SW-1000000)), Emu(350000), _['s10_vote_cn'], 28, True, WHITE, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(3050000), Emu(int(SW-1000000)), Emu(300000), _['s10_vote_en'], 18, True, GOLD, PP_ALIGN.CENTER)
cd(s, Emu(int(SW/2-2500000)), Emu(3600000), Emu(5000000), Emu(650000), GOLD)
tb(s, Emu(int(SW/2-2300000)), Emu(3680000), Emu(4600000), Emu(260000), '\u201c'+_['s10_motto']+'\u201d', 16, True, WHITE, PP_ALIGN.CENTER)
tb(s, Emu(int(SW/2-2300000)), Emu(3950000), Emu(4600000), Emu(250000), '\u201c'+_['s10_motto_en']+'\u201d', 12, False, GOLD, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(4650000), Emu(int(SW-1000000)), Emu(250000), _['s10_footer_cn'], 13, False, LTGY, PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(4880000), Emu(int(SW-1000000)), Emu(250000), _['s10_footer_en'], 11, False, MDGY, PP_ALIGN.CENTER)
print('Slide 10 done')
# === SAVE ===
prs.save(OUTPUT)
print('PPT generated successfully!')
print('Output: ' + OUTPUT)
print('Total slides: ' + str(len(prs.slides)))
